"""
Durak Makine Otomasyon — Sunum üretici.
Çıktı: cikti/durak-sunum.pptx
Çalıştırma: .venv/bin/python scripts/build_sunum.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------- palet ----------
NAVY = RGBColor(0x0F, 0x2E, 0x4F)
NAVY_DEEP = RGBColor(0x08, 0x1B, 0x30)
ACCENT = RGBColor(0xE8, 0x84, 0x2B)
ACCENT_SOFT = RGBColor(0xFD, 0xE7, 0xD0)
LIGHT = RGBColor(0xF5, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x25, 0x2E)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
GRAY_SOFT = RGBColor(0xDD, 0xE1, 0xE6)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)

FONT = "Calibri"

# ---------- yardımcılar ----------
def add_rect(slide, x, y, w, h, fill, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = fill
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_round_rect(slide, x, y, w, h, fill, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.08
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(1.25)
    s.shadow.inherit = False
    return s


def add_text(
    slide, x, y, w, h, lines, *,
    size=18, bold=False, italic=False, color=DARK, align="left", anchor="top",
    line_spacing=1.15,
):
    """lines: tek string ya da [(text, kwargs_dict), ...] biçiminde liste."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    if isinstance(lines, str):
        lines = [(lines, {})]
    for i, (text, opts) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        a = opts.get("align", align)
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }[a]
        p.line_spacing = opts.get("line_spacing", line_spacing)
        if "space_before" in opts:
            p.space_before = Pt(opts["space_before"])
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get("font", FONT)
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.italic = opts.get("italic", italic)
        run.font.color.rgb = opts.get("color", color)
    return tb


def add_line(slide, x, y, w, h, color=ACCENT, weight=2.5):
    s = slide.shapes.add_connector(1, x, y, x + w, y + h)
    s.line.color.rgb = color
    s.line.width = Pt(weight)
    return s


# ---------- ortak çatı ----------
def page_frame(slide, page_num, total=12, dark=False):
    """Tüm sayfalarda ortak: sol bant + alt bilgi şeridi."""
    bg_color = NAVY_DEEP if dark else WHITE
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, bg_color)
    # sol turuncu bant
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, ACCENT)
    # sağ alt bilgi
    foot_color = WHITE if dark else GRAY
    add_text(
        slide, Inches(0.6), SLIDE_H - Inches(0.45),
        Inches(8), Inches(0.3),
        "ESMARK Digital Marketing  ·  Durak Makine Otomasyon",
        size=10, color=foot_color,
    )
    add_text(
        slide, SLIDE_W - Inches(1.3), SLIDE_H - Inches(0.45),
        Inches(0.7), Inches(0.3),
        f"{page_num} / {total}",
        size=10, color=foot_color, align="right",
    )


def title_block(slide, kicker, title, *, dark=False):
    kicker_color = ACCENT
    title_color = WHITE if dark else NAVY
    add_text(
        slide, Inches(0.6), Inches(0.55),
        Inches(11), Inches(0.4),
        kicker, size=12, bold=True, color=kicker_color,
    )
    add_text(
        slide, Inches(0.6), Inches(0.85),
        Inches(11.5), Inches(0.9),
        title, size=32, bold=True, color=title_color,
    )
    add_line(
        slide, Inches(0.6), Inches(1.7),
        Inches(0.7), 0, color=ACCENT, weight=3,
    )


# ---------- sunum ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------- 1. Kapak ----------
def slide_kapak():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY_DEEP)
    # accent diagonal blok (sol alt)
    add_rect(s, 0, 0, Inches(0.5), SLIDE_H, ACCENT)

    add_text(
        s, Inches(1.2), Inches(1.4), Inches(11), Inches(0.6),
        "DİJİTAL YENİLEME YOL HARİTASI",
        size=14, bold=True, color=ACCENT,
    )
    add_text(
        s, Inches(1.2), Inches(2.0), Inches(11), Inches(2.0),
        "Durak Makine\nOtomasyon",
        size=64, bold=True, color=WHITE, line_spacing=1.05,
    )
    add_line(s, Inches(1.2), Inches(4.3), Inches(0.9), 0, color=ACCENT, weight=4)
    add_text(
        s, Inches(1.2), Inches(4.5), Inches(11), Inches(1.0),
        "Web sitesi · Sosyal medya · Google Business · Reklam stratejisi",
        size=20, color=WHITE,
    )

    add_text(
        s, Inches(1.2), Inches(6.6), Inches(11), Inches(0.4),
        "ESMARK Digital Marketing  ·  28 Nisan 2026",
        size=12, color=ACCENT_SOFT,
    )


# ---------- 2. Gündem ----------
def slide_gundem():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 2)
    title_block(s, "GÜNDEM", "Bugün ne konuşacağız?")

    items = [
        ("01", "Mevcut durum", "Web siteniz, Instagram'ınız ve Google Business'taki tablo"),
        ("02", "Pazar ve rakipler", "Bursa pazarındaki yeriniz, kim önde, kim niş"),
        ("03", "Sizi rakipten ayıran fark", "Söze değil, kanıta dayalı pozisyon"),
        ("04", "Yeni web sitesi planı", "Sayfa yapısı ve içerik öncelikleri"),
        ("05", "Sosyal medya ve reklam", "İçerik takvimi, hedef kitle, kampanya çerçevesi"),
    ]
    y = Inches(2.05)
    for num, head, sub in items:
        add_text(s, Inches(0.7), y, Inches(0.9), Inches(0.7),
                 num, size=34, bold=True, color=ACCENT)
        add_text(s, Inches(1.7), y + Inches(0.05), Inches(10), Inches(0.5),
                 head, size=22, bold=True, color=NAVY)
        add_text(s, Inches(1.7), y + Inches(0.55), Inches(10), Inches(0.5),
                 sub, size=14, color=GRAY)
        y += Inches(0.95)


# ---------- 3. Mevcut durum üç kanal ----------
def slide_mevcut_durum():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 3)
    title_block(s, "MEVCUT DURUM", "Üç kanalda da iş var")

    cards = [
        ("Web Sitesi", "durakmakineotomasyon.com",
         "Blog 2020'de durmuş.\nTalaşlı imalat ve büyük makine parkı sayfada hiç yok.\nAlt sayfaların bir kısmı açılmıyor.",
         RED),
        ("Instagram", "@durakmakineotomasyon",
         "Hesap var ama düzenli yönetilmiyor.\nWeb sitesinde linki bile yok.\nİçerik takvimi ve görsel arşiv eksik.",
         ACCENT),
        ("Google Business", "Yeni adres",
         "Eski adres haritada görünüyor.\nTaşınma sonrası güncelleme yapılmamış.\nMüşteri firmayı haritada bulamıyor.",
         NAVY),
    ]
    card_w = Inches(3.85)
    card_h = Inches(4.4)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.1)

    for i, (head, sub, body, color) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        # üst renk şeridi
        add_rect(s, x, y, card_w, Inches(0.5), color)
        # gövde
        add_rect(s, x, y + Inches(0.5), card_w, card_h - Inches(0.5), LIGHT)
        # başlık + alt başlık
        add_text(s, x + Inches(0.35), y + Inches(0.75), card_w - Inches(0.7), Inches(0.55),
                 head, size=22, bold=True, color=NAVY)
        add_text(s, x + Inches(0.35), y + Inches(1.3), card_w - Inches(0.7), Inches(0.4),
                 sub, size=12, color=GRAY)
        # body satırlar
        ty = y + Inches(1.85)
        for line in body.split("\n"):
            add_text(s, x + Inches(0.35), ty, card_w - Inches(0.7), Inches(0.6),
                     "•  " + line, size=13, color=DARK, line_spacing=1.25)
            ty += Inches(0.7)


# ---------- 4. Web sitesi sorunlar ----------
def slide_web_sorunlar():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 4)
    title_block(s, "WEB SİTESİ", "Tespit edilen 7 kritik eksik")

    items = [
        ("Talaşlı imalat ve büyük makine parkı sitede hiç yok",
         "Sizin en güçlü kartınız — tamamen görünmez durumda."),
        ("Blog 2020'den beri sessiz",
         "6 yıldır içerik yok. Arama motorları ve müşteriye 'durmuş firma' sinyali."),
        ("E-posta adresi yok",
         "Müşteri size sadece telefonla ulaşabiliyor. Resmi iletişim kanalı eksik."),
        ("Sosyal medya entegrasyonu yok",
         "Instagram'ınıza sitede link bile verilmemiş."),
        ("Alt sayfalar açılmıyor",
         "Bakım, revizyon, hat taşıma sayfaları ya boş ya kırık."),
        ("Foto ve video kanıt yok",
         "İşi yaptığınızı söylüyor, gösteremiyor. Rakipler bunu fark ediyor."),
        ("Yeni adresiniz hiçbir yerde yok",
         "Site, Google Business ve sosyal medya hâlâ eski lokasyonu gösteriyor."),
    ]

    y = Inches(2.05)
    for i, (head, sub) in enumerate(items):
        # numara kutusu
        add_round_rect(s, Inches(0.7), y, Inches(0.55), Inches(0.55), RED)
        add_text(s, Inches(0.7), y, Inches(0.55), Inches(0.55),
                 str(i + 1), size=18, bold=True, color=WHITE,
                 align="center", anchor="middle")
        # başlık + alt
        add_text(s, Inches(1.45), y - Inches(0.03), Inches(11), Inches(0.4),
                 head, size=15, bold=True, color=NAVY)
        add_text(s, Inches(1.45), y + Inches(0.32), Inches(11), Inches(0.4),
                 sub, size=11, color=GRAY)
        y += Inches(0.65)


# ---------- 5. Pazardaki yeriniz ----------
def slide_pazar():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 5)
    title_block(s, "PAZAR", "Bursa otomotivin başkenti — talep zaten var")

    # ana mesaj kutusu
    add_round_rect(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(1.2),
                   NAVY)
    add_text(s, Inches(1.0), Inches(2.20), Inches(11.5), Inches(0.55),
             "Bursa, Türkiye'nin otomotiv ve yan sanayi merkezi.",
             size=22, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(2.75), Inches(11.5), Inches(0.5),
             "Pres bakım, revizyon ve hat taşıma talebi sürekli — sizin işiniz organik olarak aranıyor.",
             size=14, color=ACCENT_SOFT)

    # iki sütun: rakipler / iş ortağı potansiyeli
    col_y = Inches(3.6)
    col_h = Inches(3.0)
    col_w = Inches(5.85)
    gap = Inches(0.2)

    # sol — rakipler
    add_rect(s, Inches(0.7), col_y, col_w, Inches(0.45), ACCENT)
    add_text(s, Inches(0.95), col_y + Inches(0.05), col_w, Inches(0.4),
             "DOĞRUDAN RAKİPLER", size=12, bold=True, color=WHITE)
    add_rect(s, Inches(0.7), col_y + Inches(0.45), col_w, col_h - Inches(0.45), LIGHT)

    rakipler = [
        ("Mekasis Otomasyon", "Bursa Nilüfer · ana rakip · temsilcilikler ile güçlü"),
        ("Uysal Otomasyon", "Eksantrik pres bakım, periyodik bakım odaklı"),
        ("NYT Elektronik", "Pano revizyonu, robotik pres senkronizasyonu"),
    ]
    ry = col_y + Inches(0.7)
    for name, desc in rakipler:
        add_text(s, Inches(1.0), ry, col_w - Inches(0.5), Inches(0.4),
                 name, size=15, bold=True, color=NAVY)
        add_text(s, Inches(1.0), ry + Inches(0.32), col_w - Inches(0.5), Inches(0.4),
                 desc, size=11, color=GRAY)
        ry += Inches(0.75)

    # sağ — iş ortağı potansiyeli
    rx = Inches(0.7) + col_w + gap
    add_rect(s, rx, col_y, col_w, Inches(0.45), GREEN)
    add_text(s, rx + Inches(0.25), col_y + Inches(0.05), col_w, Inches(0.4),
             "İŞ ORTAĞI POTANSİYELİ", size=12, bold=True, color=WHITE)
    add_rect(s, rx, col_y + Inches(0.45), col_w, col_h - Inches(0.45), LIGHT)

    add_text(s, rx + Inches(0.3), col_y + Inches(0.7), col_w - Inches(0.5), Inches(0.5),
             "Pres üreticileri (rakip değil, müşteri kaynağı)",
             size=14, bold=True, color=NAVY)
    uretici = "Hürsan Pres  ·  Koçlu Pres  ·  Preston  ·  PRES-SAN  ·  Meşe Makine"
    add_text(s, rx + Inches(0.3), col_y + Inches(1.15), col_w - Inches(0.5), Inches(0.6),
             uretici, size=12, color=DARK, line_spacing=1.3)
    add_text(s, rx + Inches(0.3), col_y + Inches(1.95), col_w - Inches(0.5), Inches(0.9),
             "Bu firmalarla yetkili servis veya bakım iş ortağı ilişkisi → "
             "Mekasis'in temsilcilik avantajını dengeler.",
             size=12, italic=True, color=GRAY, line_spacing=1.35)


# ---------- 6. Mekasis vs Sizler ----------
def slide_karsilastirma():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 6)
    title_block(s, "RAKİP KARŞILAŞTIRMA", "Mekasis vs Durak Makine")

    col_y = Inches(2.05)
    col_h = Inches(4.3)
    col_w = Inches(5.95)
    gap = Inches(0.2)

    # Sol — Mekasis
    add_rect(s, Inches(0.7), col_y, col_w, Inches(0.65), GRAY)
    add_text(s, Inches(0.95), col_y + Inches(0.13), col_w, Inches(0.4),
             "MEKASİS", size=11, bold=True, color=ACCENT_SOFT)
    add_text(s, Inches(0.95), col_y + Inches(0.32), col_w, Inches(0.4),
             "Söylemde önde", size=18, bold=True, color=WHITE)
    add_rect(s, Inches(0.7), col_y + Inches(0.65), col_w, col_h - Inches(0.65), LIGHT)

    mekasis_items = [
        ("+", "6 temsilcilik (Dirinler, Showa Seiki, Tecnofluid, HELM, OMPI, Vaptsarov)"),
        ("+", "'Pazar lideri' söylemi, aktif blog ve teknik içerik"),
        ("+", "Foto galeri ve sosyal medya entegrasyonu"),
        ("−", "Galeri görselleri büyük ölçüde placeholder"),
        ("−", "'5 eksen işleme' iddiası kanıtsız"),
        ("−", "Kendi makine parkı / talaşlı imalat öne çıkmıyor"),
        ("−", "Müşteri logosu ve vaka çalışması yok"),
    ]
    iy = col_y + Inches(0.85)
    for sym, txt in mekasis_items:
        sym_color = GREEN if sym == "+" else RED
        add_text(s, Inches(0.95), iy, Inches(0.3), Inches(0.4),
                 sym, size=18, bold=True, color=sym_color)
        add_text(s, Inches(1.25), iy + Inches(0.04), col_w - Inches(0.6), Inches(0.5),
                 txt, size=12, color=DARK, line_spacing=1.25)
        iy += Inches(0.45)

    # Sağ — Durak
    rx = Inches(0.7) + col_w + gap
    add_rect(s, rx, col_y, col_w, Inches(0.65), NAVY)
    add_text(s, rx + Inches(0.25), col_y + Inches(0.13), col_w, Inches(0.4),
             "DURAK MAKİNE", size=11, bold=True, color=ACCENT)
    add_text(s, rx + Inches(0.25), col_y + Inches(0.32), col_w, Inches(0.4),
             "Gerçeklikte önde", size=18, bold=True, color=WHITE)
    add_rect(s, rx, col_y + Inches(0.65), col_w, col_h - Inches(0.65), ACCENT_SOFT)

    durak_items = [
        ("✓", "Tüm işler kendi tesisinizde — taşeronsuz üretim"),
        ("✓", "Talaşlı imalat parkı: tornalar, frezeler, işleme merkezi"),
        ("✓", "Çap 780 × boy 5000 mm torna — Türkiye'de 4-5 adet"),
        ("✓", "Volan, krank mili gibi büyük parça işleme içeride"),
        ("✓", "Atölye, video ve fotoğraflarla görsel kanıt mümkün"),
        ("→", "Site bunların hiçbirini bugün gösteremiyor"),
        ("→", "Yeni site bu farkı her sayfada öne çıkaracak"),
    ]
    iy = col_y + Inches(0.85)
    for sym, txt in durak_items:
        sym_color = GREEN if sym == "✓" else ACCENT
        add_text(s, rx + Inches(0.25), iy, Inches(0.3), Inches(0.4),
                 sym, size=18, bold=True, color=sym_color)
        add_text(s, rx + Inches(0.55), iy + Inches(0.04), col_w - Inches(0.6), Inches(0.5),
                 txt, size=12, color=DARK, line_spacing=1.25)
        iy += Inches(0.45)


# ---------- 7. Sizin asıl gücünüz ----------
def slide_guc():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 7)
    title_block(s, "FARK", "Söz değil — kanıt")

    cards = [
        ("Tüm işler içeride",
         "Pres revizyonundan talaşlı imalata kadar her aşama kendi atölyenizde."),
        ("Büyük parça kapasitesi",
         "Volan, krank mili dahil ana işlemler dışarı gitmiyor — rakipler bu kapasiteyi kiralıyor."),
        ("Türkiye'de 4-5 adet",
         "Çap 780 mm × boy 5000 mm torna — bu boyut sadece sizde dahil çok az tesiste var."),
    ]
    card_w = Inches(3.85)
    card_h = Inches(2.6)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.2)
    for i, (head, body) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        add_round_rect(s, x, y, card_w, card_h, WHITE, line_color=GRAY_SOFT)
        add_round_rect(s, x, y, Inches(0.6), card_h, ACCENT, line_color=ACCENT)
        add_text(s, x + Inches(1.0), y + Inches(0.45), card_w - Inches(1.3), Inches(0.7),
                 head, size=20, bold=True, color=NAVY)
        add_text(s, x + Inches(1.0), y + Inches(1.2), card_w - Inches(1.3), Inches(1.2),
                 body, size=13, color=DARK, line_spacing=1.35)

    # alıntı
    qy = Inches(5.2)
    add_round_rect(s, Inches(1.5), qy, Inches(10.3), Inches(1.5), NAVY_DEEP)
    add_text(s, Inches(2.0), qy + Inches(0.25), Inches(9.5), Inches(0.6),
             "“ Biz burada yaptığımız her şeyi kendi tarafımızda yapıyoruz, gösterebiliriz. ”",
             size=18, italic=True, color=WHITE, line_spacing=1.3)
    add_text(s, Inches(2.0), qy + Inches(1.0), Inches(9.5), Inches(0.4),
             "— Fahrettin Durak",
             size=12, color=ACCENT)


# ---------- 8. Pazarlama mesajı ----------
def slide_mesaj():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 8, dark=True)

    add_text(
        s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.5),
        "PAZARLAMA MESAJI",
        size=14, bold=True, color=ACCENT, align="center",
    )
    add_text(
        s, Inches(1.0), Inches(2.85), Inches(11.3), Inches(2.2),
        "Sözünü kanıtlayan üretici.",
        size=64, bold=True, color=WHITE, align="center", line_spacing=1.0,
    )
    add_line(s, Inches(6.16), Inches(4.95), Inches(1.0), 0, color=ACCENT, weight=4)
    add_text(
        s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.5),
        "Dışarı iş veren değil; başkalarının yapamadığı işleri içeride yapan tesis.",
        size=20, color=ACCENT_SOFT, align="center", line_spacing=1.4,
    )


# ---------- 9. Yeni site iskeleti ----------
def slide_yeni_site():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 9)
    title_block(s, "YENİ WEB SİTESİ", "Sayfa yapısı ve öncelikler")

    # legend
    add_round_rect(s, Inches(0.7), Inches(2.0), Inches(0.18), Inches(0.18), ACCENT)
    add_text(s, Inches(0.95), Inches(1.95), Inches(3), Inches(0.3),
             "Yeni eklenen", size=11, color=GRAY)
    add_round_rect(s, Inches(2.6), Inches(2.0), Inches(0.18), Inches(0.18), NAVY)
    add_text(s, Inches(2.85), Inches(1.95), Inches(3), Inches(0.3),
             "Yenilenen", size=11, color=GRAY)

    pages = [
        ("Anasayfa", "Atölye videosu + güçlü değer önerisi", "navy"),
        ("Eksantrik Pres Revizyonu", "Vaka çalışmaları + before/after", "navy"),
        ("Hidrolik Pres", "Hizmet kapsamı + süreç adımları", "navy"),
        ("Pres Hat Taşıma", "Lojistik + montaj örnekleri", "navy"),
        ("Talaşlı İmalat", "Yeni — ana hizmet sayfası", "accent"),
        ("Makine Parkı", "Yeni — büyük torna ön plan", "accent"),
        ("Vakalar / Referanslar", "Yeni — müşteri logoları, projeler", "accent"),
        ("Blog", "Teknik içerik + SEO için yeniden açılışı", "navy"),
        ("Hakkımızda", "Kuruluş, ekip, sertifikalar", "navy"),
        ("İletişim", "Yeni adres, harita, WhatsApp CTA", "navy"),
    ]

    cols = 5
    rows = 2
    grid_w = Inches(11.9)
    grid_h = Inches(4.5)
    cell_w = Inches(2.27)
    cell_h = Inches(1.95)
    gap_x = (grid_w - cell_w * cols) / (cols - 1)
    gap_y = (grid_h - cell_h * rows) / (rows - 1)
    start_x = Inches(0.7)
    start_y = Inches(2.4)

    for i, (head, body, kind) in enumerate(pages):
        col = i % cols
        row = i // cols
        x = start_x + col * (cell_w + gap_x)
        y = start_y + row * (cell_h + gap_y)
        line_color = ACCENT if kind == "accent" else NAVY
        add_round_rect(s, x, y, cell_w, cell_h, WHITE, line_color=line_color)
        add_rect(s, x, y, cell_w, Inches(0.18), line_color)
        add_text(s, x + Inches(0.2), y + Inches(0.35), cell_w - Inches(0.4), Inches(0.7),
                 head, size=14, bold=True, color=NAVY, line_spacing=1.15)
        add_text(s, x + Inches(0.2), y + Inches(1.05), cell_w - Inches(0.4), Inches(0.85),
                 body, size=10, color=GRAY, line_spacing=1.3)


# ---------- 10. Sosyal medya & reklam ----------
def slide_sosyal_reklam():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 10)
    title_block(s, "SOSYAL MEDYA & REKLAM", "Dört kanal, tek mesaj")

    cards = [
        ("Instagram",
         "Atölye, makine ve revizyon süreci",
         ["Reels: makine zoom, before/after",
          "Story: günlük atölye anları",
          "Post: tamamlanmış işler, ekip"]),
        ("LinkedIn",
         "B2B karar vericilerine ulaşım",
         ["Hedef: satın alma + bakım müdürü",
          "Vaka çalışması paylaşımları",
          "Sektör içgörüleri ve teknik yazı"]),
        ("Google Ads",
         "Aktif aramada zirve",
         ["'eksantrik pres revizyonu Bursa'",
          "'hidrolik pres tamiri'",
          "Şehir + hizmet long-tail"]),
        ("Meta Ads",
         "Coğrafi ve sektörel hedefleme",
         ["Bursa, Kocaeli, İstanbul sanayi",
          "Otomotiv yan sanayi pozisyonları",
          "Atölye videosu ile retargeting"]),
    ]

    card_w = Inches(2.95)
    card_h = Inches(4.5)
    gap = Inches(0.13)
    total_w = card_w * 4 + gap * 3
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.1)

    for head, sub, items in cards:
        x = start_x
        add_round_rect(s, x, y, card_w, card_h, WHITE, line_color=GRAY_SOFT)
        add_rect(s, x, y, card_w, Inches(0.55), NAVY)
        add_text(s, x + Inches(0.2), y + Inches(0.13), card_w - Inches(0.4), Inches(0.4),
                 head, size=15, bold=True, color=WHITE)
        add_text(s, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(0.5),
                 sub, size=11, italic=True, color=ACCENT, line_spacing=1.25)
        iy = y + Inches(1.4)
        for it in items:
            add_text(s, x + Inches(0.2), iy, Inches(0.18), Inches(0.4),
                     "▪", size=14, bold=True, color=ACCENT)
            add_text(s, x + Inches(0.45), iy + Inches(0.02), card_w - Inches(0.65), Inches(0.85),
                     it, size=11, color=DARK, line_spacing=1.3)
            iy += Inches(0.85)
        start_x += card_w + gap


# ---------- 11. Süreç ----------
def slide_surec():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 11)
    title_block(s, "SÜREÇ", "Dört aşamada yenileme")

    steps = [
        ("01", "Keşif & İçerik",
         "Makine envanteri, atölye foto/video çekimi, mevcut müşteri ve referans listesi"),
        ("02", "Tasarım & Üretim",
         "Yeni site tasarımı, içerik yazımı, SEO yapısı, Google Business güncellemesi"),
        ("03", "Yayın & Devir",
         "Site canlıya alınır, Instagram yönetimi devralınır, ölçüm araçları kurulur"),
        ("04", "Reklam & Büyüme",
         "Meta + Google Ads kampanyaları, içerik takvimi, aylık raporlama"),
    ]

    card_w = Inches(2.95)
    card_h = Inches(4.0)
    gap = Inches(0.13)
    total_w = card_w * 4 + gap * 3
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.2)

    for i, (num, head, body) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        add_round_rect(s, x, y, card_w, card_h, WHITE, line_color=GRAY_SOFT)
        add_rect(s, x, y, card_w, Inches(0.18), ACCENT)
        add_text(s, x + Inches(0.3), y + Inches(0.4), card_w - Inches(0.6), Inches(1.0),
                 num, size=44, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.3), y + Inches(1.5), card_w - Inches(0.6), Inches(0.6),
                 head, size=18, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(2.15), card_w - Inches(0.6), Inches(1.7),
                 body, size=12, color=DARK, line_spacing=1.4)
        # ok
        if i < len(steps) - 1:
            arrow_x = x + card_w + Inches(0.005)
            add_text(s, arrow_x, y + Inches(1.7), gap, Inches(0.4),
                     "›", size=28, bold=True, color=ACCENT, align="center")

    add_text(s, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.5),
             "Toplam süre: kapsam ve içerik hazırlığına bağlı — toplantıda netleşecek.",
             size=12, italic=True, color=GRAY, align="center")


# ---------- 12. Kapanış / sonraki adımlar ----------
def slide_sonraki():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, 12, dark=True)

    add_text(s, Inches(0.7), Inches(0.55), Inches(11), Inches(0.4),
             "SONRAKİ ADIMLAR", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.9),
             "Birlikte karar vereceklerimiz",
             size=32, bold=True, color=WHITE)
    add_line(s, Inches(0.7), Inches(1.7), Inches(0.7), 0, color=ACCENT, weight=3)

    items = [
        ("01", "Yeni adres ve makine envanterinin netleştirilmesi"),
        ("02", "Foto / video çekim planı (atölye + makineler)"),
        ("03", "Web alan adı ve Instagram erişim devri"),
        ("04", "Hedef coğrafya ve müşteri profilinin onaylanması"),
        ("05", "Bütçe ve canlıya alınma tarihinin belirlenmesi"),
    ]
    y = Inches(2.1)
    for num, text in items:
        add_text(s, Inches(0.9), y, Inches(0.9), Inches(0.6),
                 num, size=26, bold=True, color=ACCENT)
        add_text(s, Inches(1.85), y + Inches(0.07), Inches(10.5), Inches(0.55),
                 text, size=18, color=WHITE)
        y += Inches(0.7)

    # alt iletişim kutusu
    add_round_rect(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.95), NAVY)
    add_text(s, Inches(1.0), Inches(6.15), Inches(11), Inches(0.4),
             "ESMARK Digital Marketing", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), Inches(6.5), Inches(11), Inches(0.4),
             "Web · Sosyal medya · Google Business · Reklam yönetimi",
             size=12, color=ACCENT_SOFT)


# ---------- üret ----------
def main():
    slide_kapak()
    slide_gundem()
    slide_mevcut_durum()
    slide_web_sorunlar()
    slide_pazar()
    slide_karsilastirma()
    slide_guc()
    slide_mesaj()
    slide_yeni_site()
    slide_sosyal_reklam()
    slide_surec()
    slide_sonraki()

    out = Path(__file__).resolve().parents[1] / "cikti" / "durak-sunum.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"OK -> {out}")


if __name__ == "__main__":
    main()
